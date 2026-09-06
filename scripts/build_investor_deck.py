"""Build the ProofyX investor pitch deck.

Structure follows the ten slides investors screen for (Introduction, Problem,
Solution, Market, Product, Traction, Team, Competition, Financials, Ask) plus a
Why-Now regulatory slide, a business-model slide and a sources appendix.

Visual language is lifted from the existing ProofyX Canva deck: deep indigo
ground, cream type, pill chips, rounded cards, the outlined hex logo and the
iridescent prism render.

Run:  python scripts/build_investor_deck.py
Out:  docs/pitch/ProofyX_Investor_Deck.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "pitch" / "assets"
OUT = ROOT / "docs" / "pitch" / "ProofyX_Investor_Deck.pptx"

# --- palette, sampled from the original deck -------------------------------
DEEP = RGBColor(0x0E, 0x04, 0x2D)
DEEP_LO = RGBColor(0x08, 0x02, 0x1C)   # gradient far corner
DEEP_HI = RGBColor(0x1A, 0x0C, 0x42)   # gradient near corner
PANEL = RGBColor(0x1C, 0x12, 0x3F)
PANEL2 = RGBColor(0x2A, 0x21, 0x45)
CREAM = RGBColor(0xFF, 0xF6, 0xED)
MUTED = RGBColor(0xA7, 0x9E, 0xC4)
TEAL = RGBColor(0x4F, 0xD1, 0xC5)
VIOLET = RGBColor(0x8B, 0x6C, 0xFF)
CYAN = RGBColor(0x4C, 0xC9, 0xF0)
MAGENTA = RGBColor(0xFF, 0x4D, 0x9D)
AMBER = RGBColor(0xFF, 0xB0, 0x20)

FONT = "Segoe UI"

SW, SH = Inches(13.333), Inches(7.5)
ML = Inches(0.62)          # left margin
CONTENT_W = SW - 2 * ML


# --- raw-XML effects python-pptx has no API for ----------------------------
def _sub(parent, tag, **attrs):
    el = etree.SubElement(parent, qn(tag))
    for k, v in attrs.items():
        el.set(k, str(v))
    return el


def soft_fill(shape, colour, alpha_pct, blur_in=0.0):
    """Solid fill at partial opacity, optionally with a soft (blurred) edge.

    Used for the ambient colour washes behind content. python-pptx exposes
    neither alpha nor softEdge, so both are written into the shape's spPr.
    """
    spPr = shape._element.spPr
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill"):
        existing = spPr.find(qn(tag))
        if existing is not None:
            spPr.remove(existing)

    fill = etree.SubElement(spPr, qn("a:solidFill"))
    clr = _sub(fill, "a:srgbClr", val=f"{colour[0]:02X}{colour[1]:02X}{colour[2]:02X}")
    _sub(clr, "a:alpha", val=int(alpha_pct * 1000))
    # solidFill must precede ln/effectLst in spPr's schema order.
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
        spPr.append(ln)
    if blur_in:
        effects = etree.SubElement(spPr, qn("a:effectLst"))
        _sub(effects, "a:softEdge", rad=int(Inches(blur_in)))
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def glow(slide, cx, cy, diameter, colour, alpha_pct=14, blur_in=0.9):
    """Ambient light bloom, centred on (cx, cy). Purely decorative."""
    d = Inches(diameter)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, int(cx - d / 2), int(cy - d / 2), d, d)
    return soft_fill(shape, colour, alpha_pct, blur_in)


# --- primitives ------------------------------------------------------------
def new_slide(prs, bloom=True):
    """Gradient ground plus an optional violet/teal bloom, so no slide is a
    flat rectangle of colour behind the content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.gradient()
    bg.fill.gradient_angle = 315.0
    stops = bg.fill.gradient_stops
    stops[0].color.rgb = DEEP_HI
    stops[0].position = 0.0
    stops[1].color.rgb = DEEP_LO
    stops[1].position = 1.0
    bg.line.fill.background()
    bg.shadow.inherit = False

    if bloom:
        glow(slide, Inches(12.4), Inches(0.1), 8.6, VIOLET, alpha_pct=11,
             blur_in=2.2)
        glow(slide, Inches(-0.2), Inches(7.7), 7.0, TEAL, alpha_pct=7,
             blur_in=2.0)
    return slide


def textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, line=1.15, space_after=0):
    """runs: list of (text, size_pt, colour, bold) or (text, size, colour, bold, spacing).

    A non-positive extent produces a file PowerPoint refuses to open, so a bad
    layout calculation fails loudly here rather than silently at render time.
    """
    if w <= 0 or h <= 0:
        raise ValueError(
            f"textbox needs a positive extent, got w={Emu(int(w)).inches:.2f}in "
            f"h={Emu(int(h)).inches:.2f}in for {runs[0][0][:40]!r}")
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, run in enumerate(runs):
        text, size, colour, bold = run[:4]
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = run[4] if len(run) > 4 else line
        para.space_after = Pt(space_after)
        r = para.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = colour
        r.font.bold = bold
        r.font.name = FONT
    return box


def card(slide, x, y, w, h, fill=None, outline=CREAM, radius=0.06, line_w=1.1,
         gradient=True):
    """Rounded surface. Filled cards get a top-lit gradient by default, which
    stops a grid of them reading as flat swatches."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = radius
    if fill is None:
        shape.fill.background()
    elif gradient:
        shape.fill.gradient()
        shape.fill.gradient_angle = 270.0
        stops = shape.fill.gradient_stops
        stops[0].color.rgb = RGBColor(min(fill[0] + 14, 255),
                                      min(fill[1] + 12, 255),
                                      min(fill[2] + 20, 255))
        stops[0].position = 0.0
        stops[1].color.rgb = fill
        stops[1].position = 1.0
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    shape.text_frame.text = ""
    return shape


def pill(slide, x, y, text, w=None, fill=None, outline=CREAM, colour=CREAM,
         size=11.5):
    w = w or Inches(0.20 + 0.095 * len(text))
    h = Inches(0.36)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.5
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = outline
    shape.line.width = Pt(1.0)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = colour
    r.font.name = FONT
    return shape


def rule(slide, x, y, w, colour=TEAL, thickness=3.0):
    """Short accent bar. Anchors a headline without adding another box."""
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w,
                                 Pt(thickness))
    bar.adjustments[0] = 0.5
    bar.fill.solid()
    bar.fill.fore_color.rgb = colour
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def chrome(slide, section):
    """Logo, section chip and accent rule — the furniture on content slides."""
    slide.shapes.add_picture(str(ASSETS / "logo_white.png"),
                             Inches(0.42), Inches(0.24), height=Inches(0.86))
    pill(slide, Inches(2.55), Inches(0.40), section.upper(),
         fill=None, outline=VIOLET, colour=CREAM, size=10.5)


def heading(slide, light, bold=None, y=1.02, size=40):
    """Two-tone headline matching the original deck's 'Market Analysis' style."""
    box = slide.shapes.add_textbox(ML, Inches(y), CONTENT_W, Inches(0.95))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = light
    r1.font.size = Pt(size)
    r1.font.bold = True
    r1.font.color.rgb = CREAM
    r1.font.name = FONT
    if bold:
        r2 = p.add_run()
        r2.text = " " + bold
        r2.font.size = Pt(size)
        r2.font.bold = True
        r2.font.color.rgb = TEAL
        r2.font.name = FONT
    return box


def kicker(slide, text, y=1.86, colour=MUTED, size=13.5, w=None, bar=True):
    if bar:
        rule(slide, ML, Inches(y - 0.19), Inches(1.05))
    return textbox(slide, ML, Inches(y), w or CONTENT_W, Inches(0.5),
                   [(text, size, colour, False)], line=1.3)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def stat_tile(slide, x, y, w, h, value, label, accent=TEAL, vsize=30, lsize=11.5,
              side_by_side=False):
    """Vertical by default; side_by_side puts the number left of its label."""
    card(slide, x, y, w, h, fill=PANEL, outline=accent, line_w=1.2)
    pad = Inches(0.24)
    if side_by_side:
        num_w = Inches(1.72)
        textbox(slide, x + pad, y, num_w, h,
                [(value, vsize, accent, True)], anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, x + pad + num_w, y, w - pad - num_w - Inches(0.20), h,
                [(label, lsize, CREAM, False)], line=1.28,
                anchor=MSO_ANCHOR.MIDDLE)
    else:
        # One line of the number at ~1.35 line-height, converted pt -> inches.
        value_h = Inches(vsize * 1.35 / 72.0)
        textbox(slide, x + pad, y + Inches(0.16), w - 2 * pad, value_h,
                [(value, vsize, accent, True)])
        label_y = y + Inches(0.16) + value_h
        textbox(slide, x + pad, label_y, w - 2 * pad,
                max(h - (label_y - y) - Inches(0.12), Inches(0.30)),
                [(label, lsize, CREAM, False)], line=1.25)


def bullet_card(slide, x, y, w, h, tag, title, body, accent=CYAN,
                fill=PANEL, tsize=16, bsize=11, title_lines=1):
    """title_lines reserves vertical room so a wrapping title never sits on the body."""
    card(slide, x, y, w, h, fill=fill, outline=accent, line_w=1.2)
    pad = Inches(0.26)
    title_h = Inches(0.30 * title_lines * (tsize / 15.0))
    body_y = y + Inches(0.52) + title_h + Inches(0.14)
    textbox(slide, x + pad, y + Inches(0.20), w - 2 * pad, Inches(0.28),
            [(tag.upper(), 9.5, accent, True)])
    textbox(slide, x + pad, y + Inches(0.52), w - 2 * pad, title_h,
            [(title, tsize, CREAM, True)], line=1.1)
    textbox(slide, x + pad, body_y, w - 2 * pad,
            h - (body_y - y) - Inches(0.20),
            [(body, bsize, MUTED, False)], line=1.34)


def table(slide, x, y, w, rows, col_w, header_fill=PANEL2,
          row_h=Inches(0.42), header_h=Inches(0.46), fsize=10.5,
          highlight_row=None, highlight=TEAL):
    n_rows, n_cols = len(rows), len(rows[0])
    h = header_h + row_h * (n_rows - 1)
    gf = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = cw
    tbl.rows[0].height = header_h
    for r in range(1, n_rows):
        tbl.rows[r].height = row_h

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            cell.margin_left = Inches(0.13)
            cell.margin_right = Inches(0.10)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = header_fill
            elif highlight_row is not None and r == highlight_row:
                cell.fill.fore_color.rgb = RGBColor(0x1B, 0x3B, 0x44)
            else:
                cell.fill.fore_color.rgb = PANEL
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.name = FONT
            run.font.size = Pt(fsize)
            if r == 0:
                run.font.bold = True
                run.font.color.rgb = MUTED
            elif highlight_row is not None and r == highlight_row:
                run.font.bold = True
                run.font.color.rgb = highlight
            else:
                run.font.color.rgb = CREAM
    return gf


def picture_fit(slide, name, x, y, max_w, max_h):
    """Place an image scaled to fit inside the box, centred."""
    from PIL import Image
    path = ASSETS / name
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    return slide.shapes.add_picture(str(path),
                                    x + int((max_w - w) / 2),
                                    y + int((max_h - h) / 2),
                                    width=w, height=h)


# --- slides ----------------------------------------------------------------
def s01_title(prs):
    s = new_slide(prs)
    s.shapes.add_picture(str(ASSETS / "logo_white.png"),
                         Inches(0.42), Inches(0.24), height=Inches(1.30))

    textbox(s, ML, Inches(1.80), Inches(11.0), Inches(1.5),
            [("ProofyX", 92, CREAM, True)])
    textbox(s, ML + Inches(0.06), Inches(3.42), Inches(6.6), Inches(1.7),
            [("Is this real?", 30, TEAL, True),
             ("ProofyX answers that in seconds — for any image, video or voice — "
              "and shows the evidence behind the answer.", 17, CREAM, False)],
            line=1.34, space_after=10)

    s.shapes.add_picture(str(ASSETS / "prism.png"),
                         Inches(7.05), Inches(3.85), width=Inches(6.05))

    pill(s, ML, Inches(6.52), "Investor Pitch Deck", w=Inches(2.05))
    pill(s, ML + Inches(2.20), Inches(6.52), "Seed Round", w=Inches(1.55),
         fill=CREAM, outline=CREAM, colour=DEEP)
    notes(s, "15s. 'ProofyX answers one question: is this real? For any image, "
             "video or voice file, in seconds, with the evidence attached.' "
             "Do not explain the architecture yet.")
    return s


def s02_problem(prs):
    s = new_slide(prs)
    chrome(s, "01 / Problem")
    heading(s, "Nobody can tell", "what is real any more.")
    kicker(s, "Synthetic media crossed the threshold where human review stops "
              "working — and the losses are already being booked.", w=Inches(7.4))

    picture_fit(s, "fig_detection_gap.png", Inches(6.95), Inches(2.55),
                Emu(int(Inches(5.75))), Emu(int(Inches(3.05))))

    tiles = [
        ("8M+", "Deepfake files in circulation in 2025,\nup from 500K in 2023", MAGENTA),
        ("$893M", "AI-enabled fraud losses logged by the FBI\nin 2025, across 22,364 complaints", MAGENTA),
        ("$280K", "Average loss per single\ndeepfake fraud incident", AMBER),
    ]
    y = Inches(2.52)
    for value, label, accent in tiles:
        stat_tile(s, ML, y, Inches(6.05), Inches(1.14), value, label,
                  accent=accent, vsize=27, lsize=11, side_by_side=True)
        y += Inches(1.26)

    textbox(s, ML, Inches(6.42), Inches(6.05), Inches(0.9),
            [("Detection is now a compliance problem, not a curiosity. "
              "Banks, newsrooms and courts need an answer they can put in a file.",
              11, CREAM, False)], line=1.3)
    textbox(s, Inches(6.95), Inches(5.86), Inches(5.75), Inches(0.6),
            [("iProov tested 2,000 UK and US consumers on real and fake media. "
              "Six in ten backed themselves; one in a thousand got every item "
              "right.", 10, MUTED, False)], line=1.28)
    notes(s, "35s. Lead with the confidence gap on the right: people are sure "
             "they can tell, and they cannot. Then the money. Land on: this is "
             "already a booked "
             "loss line, not a future risk.")
    return s


def s03_why_now(prs):
    s = new_slide(prs)
    chrome(s, "02 / Why now")
    heading(s, "Detection just became", "legally mandatory.")
    kicker(s, "Regulation is converting deepfake detection from a nice-to-have "
              "into compulsory spend — on a clock we can name.")

    items = [
        ("EU AI ACT", "Article 50",
         "From 2 August 2026, AI-generated content must be machine-readable and "
         "detectable. Article 99 penalties reach EUR 15M or 3% of global turnover.",
         VIOLET),
        ("UNITED STATES", "TAKE IT DOWN Act + 46 states",
         "Signed May 2025. 169 state deepfake laws since 2022, 146 bills in 2025 "
         "alone. FinCEN warned financial institutions in Q4 2025.", CYAN),
        ("INDIA", "IT Rules amendment, 2026",
         "Synthetic-media labelling and takedown duties for intermediaries — the "
         "market where our first design partners sit.", TEAL),
        ("GLOBAL", "China, South Korea, FATF",
         "FATF named deepfakes a direct threat to AML and customer due diligence. "
         "South Korea legislated after a 1,625% spike in deepfake fraud.", AMBER),
    ]
    gap = Inches(0.24)
    w = (CONTENT_W - 3 * gap) / 4
    x = ML
    for tag, title, body, accent in items:
        bullet_card(s, x, Inches(2.50), w, Inches(3.05), tag, title, body,
                    accent=accent, tsize=14.5, bsize=11.5, title_lines=2)
        x += w + gap

    card(s, ML, Inches(5.80), CONTENT_W, Inches(0.86), fill=PANEL,
         outline=TEAL, line_w=1.4)
    textbox(s, ML + Inches(0.30), Inches(6.02), CONTENT_W - Inches(0.60),
            Inches(0.5),
            [("The buyer no longer has to be convinced the problem is real — "
              "a regulator already told them. That changes the sales cycle from "
              "evangelism to procurement.", 13, CREAM, True)], line=1.25)
    notes(s, "30s. This is the slide that answers 'why now'. The EU date is "
             "fixed and close. Budget already exists on the other side of it.")
    return s


def s04_solution(prs):
    s = new_slide(prs)
    chrome(s, "03 / Solution")
    heading(s, "One number, and", "the reason behind it.")
    kicker(s, "ProofyX returns a 0-100 Trust Score for any image, video or audio "
              "file — plus the heatmap and per-model evidence that produced it.",
           w=Inches(7.5))

    # Sized by height so the band cannot run off the bottom of the slide.
    picture_fit(s, "modality_flow.png", ML, Inches(5.34),
                Emu(int(CONTENT_W)), Emu(int(Inches(1.82))))

    items = [
        ("EVIDENCE, NOT A VERDICT", "Explainable by construction",
         "GradCAM heatmaps show which pixels drove the score, and every model in "
         "the ensemble reports separately. Competitors return a binary label; a "
         "compliance officer cannot file a binary label.", TEAL),
        ("BREADTH", "Image, video, audio and cross-modal",
         "One API covers all three modalities and scores a video on its frames "
         "AND its soundtrack, surfacing disagreement between them rather than "
         "averaging it away.", CYAN),
        ("DURABILITY", "Modular against the arms race",
         "Detection is a moving target. Models are swappable behind a stable "
         "scoring interface, so a new generator is a model update, not a rewrite.",
         VIOLET),
    ]
    gap = Inches(0.28)
    w = (CONTENT_W - 2 * gap) / 3
    x = ML
    for tag, title, body, accent in items:
        bullet_card(s, x, Inches(2.52), w, Inches(2.62), tag, title, body,
                    accent=accent, tsize=15, bsize=11.5, title_lines=2)
        x += w + gap
    notes(s, "30s. The unique value proposition is the second half of the "
             "sentence: everyone can say fake or real. We say why, in a form "
             "that survives an audit.")
    return s


def s05_product(prs):
    s = new_slide(prs)
    chrome(s, "04 / Product")
    heading(s, "Upload. Score.", "Act.")
    kicker(s, "Live product: React dashboard, documented REST API, ten models "
              "loaded, running today on image, video and audio.", w=Inches(6.4))

    pic_w, pic_h = Inches(6.05), Inches(3.35)
    s.shapes.add_picture(str(ASSETS / "ui_dashboard.png"),
                         Inches(6.66), Inches(2.44), width=pic_w, height=pic_h)
    card(s, Inches(6.66), Inches(2.44), pic_w, pic_h, fill=None,
         outline=VIOLET, line_w=1.4)
    textbox(s, Inches(6.66), Inches(5.88), pic_w, Inches(0.4),
            [("Concept render of the analysis view. Replace with a live "
              "screenshot before pitching — the product is running.",
              10.5, MUTED, False)], line=1.25)

    steps = [
        ("1", "Upload", "Drag in a file, or POST to /api/v1. JPG PNG WebP, "
                        "MP4 MOV MKV, WAV MP3 FLAC."),
        ("2", "Analyse", "Ten models score in parallel; learned fusion "
                         "reconciles them with calibrated confidence."),
        ("3", "Explain", "Trust Score 0-100, GradCAM heatmap, per-model "
                         "scores and an EXIF forensics report."),
        ("4", "Act", "Four-tier verdict, exportable report, and a "
                     "one-click cybercrime complaint filing in India."),
    ]
    y = Inches(2.44)
    for num, title, body in steps:
        card(s, ML, y, Inches(5.80), Inches(0.78), fill=PANEL,
             outline=PANEL2, line_w=1.0)
        c = card(s, ML + Inches(0.16), y + Inches(0.17), Inches(0.44),
                 Inches(0.44), fill=TEAL, outline=None, radius=0.5)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = DEEP
        r.font.name = FONT
        textbox(s, ML + Inches(0.74), y + Inches(0.12), Inches(1.30),
                Inches(0.34), [(title, 13.5, CREAM, True)])
        textbox(s, ML + Inches(1.98), y + Inches(0.11), Inches(3.68),
                Inches(0.58), [(body, 11, MUTED, False)], line=1.26)
        y += Inches(0.88)

    verdicts = [("LIKELY MANIPULATED", MAGENTA), ("POSSIBLY MANIPULATED", AMBER),
                ("UNCERTAIN", MUTED), ("LIKELY AUTHENTIC", TEAL)]
    x = ML
    for label, colour in verdicts:
        w = Inches(0.30 + 0.083 * len(label))
        pill(s, x, Inches(6.10), label, w=w, fill=None, outline=colour,
             colour=colour, size=10)
        x += w + Inches(0.12)
    textbox(s, ML, Inches(6.62), Inches(5.8), Inches(0.4),
            [("Four-tier verdict — we say 'uncertain' when we are, rather than "
              "forcing a coin-flip.", 10.5, MUTED, False)], line=1.25)
    notes(s, "40s. If there is a laptop in the room, demo instead of showing "
             "this: upload one fake, land on the heatmap. 60 seconds, one flow.")
    return s


def s06_traction(prs):
    s = new_slide(prs)
    chrome(s, "05 / Validation")
    heading(s, "What is measured,", "and what is not.")
    kicker(s, "Every number below is reproducible from this repository. We label "
              "targets as targets.")

    rows = [
        ["Capability", "Status", "Measured result", "Evidence"],
        ["CorefakeNet unified image model", "Shipped",
         "82.5% acc · 90.9% ROC-AUC · F1 81.3", "332 held-out samples"],
        ["Single-image latency (CPU)", "Shipped",
         "564 ms — 4.9x faster than the ensemble", "Benchmark run"],
        ["Video pipeline, end to end", "Shipped",
         "Verified at 1/2/10/20 FPS, all 200 OK", "docs section 18"],
        ["Audio deepfake detection", "Shipped",
         "Wav2Vec2 backbone — 97.9% on the model card", "Upstream benchmark"],
        ["REST API + React dashboard", "Live",
         "10 models loaded, OpenAPI documented, rate-limited", "Deployed build"],
        ["7-model ensemble accuracy", "In progress",
         "95%+ target — fusion calibration being corrected", "Not yet claimable"],
        ["Paying customers", "Not yet", "Zero. Seeking design partners.",
         "Honest baseline"],
    ]
    col_w = [Inches(3.95), Inches(1.55), Inches(4.45), Inches(2.15)]
    table(s, ML, Inches(2.50), CONTENT_W, rows, col_w, fsize=10,
          row_h=Inches(0.46), header_h=Inches(0.44))

    card(s, ML, Inches(6.28), CONTENT_W, Inches(0.78), fill=PANEL,
         outline=AMBER, line_w=1.3)
    textbox(s, ML + Inches(0.30), Inches(6.48), CONTENT_W - Inches(0.60),
            Inches(0.44),
            [("We are pre-revenue. The ask funds the move from a working, "
              "measured prototype to a calibrated product with design partners "
              "under contract.", 12, CREAM, True)], line=1.25)
    notes(s, "35s. Do not skip the last two rows. Volunteering the weak number "
             "is what makes the strong ones believable in diligence.")
    return s


def s07_market(prs):
    s = new_slide(prs)
    chrome(s, "06 / Market")
    heading(s, "A $170M market", "compounding at 47.6%.")
    kicker(s, "Early, mandated, and accelerating fast.", w=Inches(6.2))

    picture_fit(s, "fig_market_growth.png", Inches(6.50), Inches(2.28),
                Emu(int(Inches(6.22))), Emu(int(Inches(4.10))))

    tiles = [
        ("$5.6B", "Deepfake detection market by 2034\n(from $170M in 2025)", CYAN),
        ("66.7%", "Of the market is video and image\ndetection — our core strength", VIOLET),
        ("49.2%", "Media and journalism share, with\nfinancial services fastest-growing", TEAL),
    ]
    y = Inches(2.40)
    for value, label, accent in tiles:
        stat_tile(s, ML, y, Inches(5.60), Inches(1.26), value, label,
                  accent=accent, vsize=28)
        y += Inches(1.36)

    textbox(s, ML, Inches(6.52), Inches(5.60), Inches(0.9),
            [("Adjacent budget lines we can reach from the same API: digital "
              "identity verification (~$18B by 2027) and content moderation "
              "(~$12B by 2027).", 11.5, MUTED, False)], line=1.3)
    notes(s, "25s. The point is the CAGR plus the mandate on the previous "
             "slide. A 47.6% CAGR with a legal forcing function is rare.")
    return s


def s08_opportunity(prs):
    s = new_slide(prs)
    chrome(s, "06 / Market")
    heading(s, "TAM, SAM and the", "share we are underwriting.")
    kicker(s, "Sized bottom-up from the segments a detection API can actually "
              "invoice, not from the whole cybersecurity market.")

    picture_fit(s, "fig_tam_sam_som.png", ML, Inches(2.40),
                Emu(int(Inches(7.35))), Emu(int(Inches(4.05))))

    x = Inches(8.30)
    w = Inches(4.40)
    textbox(s, x, Inches(2.46), w, Inches(0.4),
            [("WHERE THE FIRST DOLLARS COME FROM", 10, VIOLET, True)])
    segs = [
        ("Financial services", "KYC/AML and voice-auth fraud. $0.50-1.50 per "
                               "verification. Fastest-growing vertical."),
        ("Media and newsrooms", "Pre-publication verification. Largest current "
                                "share of the market at ~49%."),
        ("Government and legal", "Evidence authentication and counter-"
                                 "disinformation. $0.82B defence segment in 2025."),
    ]
    y = Inches(2.88)
    for title, body in segs:
        card(s, x, y, w, Inches(1.02), fill=PANEL, outline=PANEL2, line_w=1.0)
        textbox(s, x + Inches(0.24), y + Inches(0.15), w - Inches(0.48),
                Inches(0.3), [(title, 13, TEAL, True)])
        textbox(s, x + Inches(0.24), y + Inches(0.46), w - Inches(0.48),
                Inches(0.5), [(body, 11, MUTED, False)], line=1.28)
        y += Inches(1.14)

    textbox(s, x, Inches(6.34), w, Inches(0.6),
            [("Comparable: Doppel grew enterprise customers 400% in one year.",
              11, MUTED, False)], line=1.28)
    notes(s, "25s. Be ready for 'how did you get to SOM' — it is 2-4% of SAM "
             "at year 5, and the assumptions are in the appendix.")
    return s


def s09_competition(prs):
    s = new_slide(prs)
    chrome(s, "07 / Competition")
    heading(s, "A funded field —", "and a gap in it.")
    kicker(s, "26 funded companies have raised $259M between them. None of the "
              "leaders ship breadth and explainability in the same product.")

    rows = [
        ["Player", "Raised", "Modalities", "Explainable", "Fast mode"],
        ["Reality Defender", "$52.4M", "Multi", "Partial", "No"],
        ["Doppel", "$124M", "Social eng.", "No", "Real-time"],
        ["Pindrop", "Major VC", "Voice only", "No", "Yes"],
        ["GetReal Security", "$17.5M", "Visual", "Forensic", "No"],
        ["Sensity AI", "Undisclosed", "Visual", "Forensic", "No"],
        ["Intel FakeCatcher", "Internal", "Video only", "Physiological", "Yes"],
        ["ProofyX", "Pre-seed", "Image+Video+Audio", "GradCAM + per-model", "Yes (4.9x)"],
    ]
    col_w = [Inches(2.60), Inches(1.55), Inches(2.35), Inches(2.85), Inches(1.75)]
    table(s, ML, Inches(2.48), Inches(11.10), rows, col_w, fsize=10.5,
          row_h=Inches(0.40), header_h=Inches(0.44), highlight_row=7)

    card(s, ML, Inches(5.86), CONTENT_W, Inches(1.06), fill=PANEL,
         outline=TEAL, line_w=1.4)
    textbox(s, ML + Inches(0.30), Inches(6.04), CONTENT_W - Inches(0.60),
            Inches(0.8),
            [("Our narrow, defensible claim:", 12, TEAL, True),
             ("ProofyX is the only stack combining all three modalities, "
              "per-model explainability and a fast single-pass mode behind one "
              "API — the shape a regulated buyer needs to file a decision, not "
              "just make one. We are not claiming the best single-model accuracy; "
              "we are claiming the most auditable answer.", 11.5, CREAM, False)],
            line=1.3, space_after=3)
    notes(s, "30s. Never say 'no competition'. Name the funded players first, "
             "then the one row where we differ. Expect a question on moat — "
             "answer with the explainability workflow and the modular retrain "
             "loop, not the model weights.")
    return s


def s10_business_model(prs):
    s = new_slide(prs)
    chrome(s, "08 / Business model")
    heading(s, "Four ways to charge", "for the same API call.")
    kicker(s, "Land through self-serve and the API, expand into enterprise "
              "licences. Benchmarked against what this market already pays.")

    rows = [
        ["Stream", "Price", "Target customer", "Why they buy"],
        ["API-as-a-service", "$0.25-0.50 / analysis", "Platforms, integrators",
         "Volume moderation, drop-in call"],
        ["SaaS subscription", "$99-499 / month", "Newsrooms, SMBs",
         "Desk-level verification workflow"],
        ["Per-check KYC", "$0.50-1.50 / check", "Financial services",
         "Regulatory duty, priced per user"],
        ["Enterprise licence", "$50K-500K / year", "Banks, government",
         "On-prem, SLA, audit trail"],
    ]
    col_w = [Inches(2.60), Inches(2.45), Inches(3.05), Inches(3.99)]
    table(s, ML, Inches(2.48), CONTENT_W, rows, col_w, fsize=10.5,
          row_h=Inches(0.52), header_h=Inches(0.46))

    y = Inches(5.44)
    gap = Inches(0.24)
    w = (CONTENT_W - 2 * gap) / 3
    facts = [
        ("Benchmark", "Sightengine charges $99/mo for 40,000 images; "
                      "deepidv $19/mo + $0.50/check."),
        ("Gross margin", "CPU inference today; GPU batch on the roadmap moves "
                         "unit cost down as volume rises."),
        ("Expansion", "Same API serves moderation, KYC and evidence — one "
                      "integration, three budget lines."),
    ]
    x = ML
    for title, body in facts:
        card(s, x, y, w, Inches(1.16), fill=PANEL, outline=PANEL2, line_w=1.0)
        textbox(s, x + Inches(0.24), y + Inches(0.16), w - Inches(0.48),
                Inches(0.3), [(title.upper(), 10.5, CYAN, True)])
        textbox(s, x + Inches(0.24), y + Inches(0.48), w - Inches(0.48),
                Inches(0.6), [(body, 11, MUTED, False)], line=1.3)
        x += w + gap
    notes(s, "20s. Investors want to know the land-and-expand path. Self-serve "
             "API is the wedge, enterprise licence is the value.")
    return s


def s11_financials(prs):
    s = new_slide(prs)
    chrome(s, "09 / Financials")
    heading(s, "Five-year path", "to $100-200M ARR.")
    kicker(s, "Conservative and target cases. Pre-revenue today — these are "
              "projections, and the assumptions are listed opposite.",
           w=Inches(6.4))

    picture_fit(s, "fig_revenue_projection.png", ML, Inches(2.28),
                Emu(int(Inches(7.05))), Emu(int(Inches(4.10))))

    x = Inches(8.05)
    w = Inches(4.65)
    textbox(s, x, Inches(2.44), w, Inches(0.4),
            [("ASSUMPTIONS BEHIND THE CURVE", 10, VIOLET, True)])
    assumptions = [
        ("Year 1", "3-6 design partners at $50-150K. No self-serve revenue "
                   "assumed."),
        ("Year 2-3", "Self-serve API opens; enterprise ACV rises as the EU AI "
                     "Act deadline bites."),
        ("Year 4-5", "2-4% of a $5-8B SAM. Requires GPU inference and a "
                     "compliance-certified deployment."),
        ("Cost base", "Team and compute dominate; inference cost per call falls "
                      "with batch GPU."),
    ]
    y = Inches(2.86)
    for title, body in assumptions:
        card(s, x, y, w, Inches(0.86), fill=PANEL, outline=PANEL2, line_w=1.0)
        textbox(s, x + Inches(0.22), y + Inches(0.13), w - Inches(0.44),
                Inches(0.26), [(title, 11.5, TEAL, True)])
        textbox(s, x + Inches(0.22), y + Inches(0.40), w - Inches(0.44),
                Inches(0.44), [(body, 11, MUTED, False)], line=1.26)
        y += Inches(0.96)
    notes(s, "25s. Say the word 'projection' out loud. If asked for the model, "
             "the spreadsheet is in the data room — do not improvise numbers.")
    return s


def s12_team(prs):
    s = new_slide(prs)
    chrome(s, "10 / Team")
    heading(s, "Builders who already", "shipped the hard part.")
    kicker(s, "Ten models, three modalities and a live API — built by this team "
              "before raising a rupee.")

    members = [
        ("Muhammed Sayeedur Rahman S", "Founder / ML & Platform",
         "Owns the detection stack and API: model training, fusion calibration, "
         "FastAPI backend and the compliance/traceability layer."),
        ("Mohamed Abrar", "Co-founder / Product & Frontend",
         "Largest contributor by commit count. Owns the React dashboard, the "
         "analysis workflows and the deployment pipeline."),
        ("Mohamed Shahid", "Co-founder / Infrastructure & QA",
         "Owns CI/CD, containerisation, security scanning and the test suite "
         "across the backend."),
    ]
    gap = Inches(0.30)
    w = (CONTENT_W - 2 * gap) / 3
    x = ML
    for name, role, body in members:
        card(s, x, Inches(2.50), w, Inches(2.80), fill=PANEL, outline=CYAN,
             line_w=1.2)
        c = card(s, x + Inches(0.28), Inches(2.72), Inches(0.60), Inches(0.60),
                 fill=PANEL2, outline=CYAN, radius=0.5)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "".join(w2[0] for w2 in name.split()[:2]).upper()
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = CYAN
        r.font.name = FONT
        textbox(s, x + Inches(0.28), Inches(3.40), w - Inches(0.56),
                Inches(0.58), [(name, 14, CREAM, True)], line=1.15)
        textbox(s, x + Inches(0.28), Inches(4.02), w - Inches(0.56),
                Inches(0.28), [(role, 10.5, TEAL, True)])
        textbox(s, x + Inches(0.28), Inches(4.34), w - Inches(0.56),
                Inches(0.84), [(body, 11, MUTED, False)], line=1.3)
        x += w + gap

    card(s, ML, Inches(5.42), CONTENT_W, Inches(1.16), fill=PANEL,
         outline=AMBER, line_w=1.3)
    textbox(s, ML + Inches(0.30), Inches(5.60), CONTENT_W - Inches(0.60),
            Inches(0.86),
            [("What we are missing, and are hiring for", 12, AMBER, True),
             ("A commercial lead with enterprise security or regtech sales "
              "experience, and a senior advisor from media forensics or "
              "financial-crime compliance. Both are line items in the ask.",
              11, CREAM, False)], line=1.28, space_after=3)
    notes(s, "25s. Name the gap yourself. An investor who spots a missing "
             "hustler before you name it will discount everything else.")
    return s


def s13_ask(prs):
    s = new_slide(prs)
    chrome(s, "11 / The ask")
    heading(s, "Raising", "$1.5M-$2.5M seed.")
    kicker(s, "18-24 months of runway to reach calibrated accuracy, paying "
              "design partners and an EU AI Act-ready deployment.")

    uses = [
        ("40%", "Engineering", "Two ML engineers and one platform engineer. "
                               "Fix fusion calibration, ship GPU inference.", TEAL),
        ("25%", "Go-to-market", "Commercial lead plus design-partner pilots in "
                                "financial services and media.", CYAN),
        ("20%", "Compute & data", "GPU training and inference, licensed "
                                  "evaluation datasets, red-team generation.", VIOLET),
        ("15%", "Compliance", "Security certification, audit trail and the "
                              "evidence pack regulated buyers require.", AMBER),
    ]
    gap = Inches(0.24)
    w = (CONTENT_W - 3 * gap) / 4
    x = ML
    for pct, title, body, accent in uses:
        card(s, x, Inches(2.48), w, Inches(2.06), fill=PANEL, outline=accent,
             line_w=1.2)
        textbox(s, x + Inches(0.24), Inches(2.66), w - Inches(0.48),
                Inches(0.6), [(pct, 30, accent, True)])
        textbox(s, x + Inches(0.24), Inches(3.24), w - Inches(0.48),
                Inches(0.3), [(title, 13, CREAM, True)])
        textbox(s, x + Inches(0.24), Inches(3.58), w - Inches(0.48),
                Inches(0.8), [(body, 11, MUTED, False)], line=1.3)
        x += w + gap

    textbox(s, ML, Inches(4.78), CONTENT_W, Inches(0.4),
            [("WHAT THE MONEY BUYS — MILESTONES WE WILL BE JUDGED ON",
              10, VIOLET, True)])
    milestones = [
        ("Month 6", "Ensemble calibration fixed; published accuracy on a public "
                    "benchmark, third-party reproducible."),
        ("Month 12", "3-6 design partners live under paid pilots; first "
                     "enterprise contract signed."),
        ("Month 18", "EU AI Act-ready deployment, SOC 2 in progress, GPU "
                     "inference cutting unit cost."),
    ]
    gap2 = Inches(0.24)
    w2 = (CONTENT_W - 2 * gap2) / 3
    x = ML
    for title, body in milestones:
        card(s, x, Inches(5.20), w2, Inches(1.22), fill=PANEL, outline=PANEL2,
             line_w=1.0)
        textbox(s, x + Inches(0.24), Inches(5.36), w2 - Inches(0.48),
                Inches(0.3), [(title, 12.5, TEAL, True)])
        textbox(s, x + Inches(0.24), Inches(5.70), w2 - Inches(0.48),
                Inches(0.66), [(body, 11, MUTED, False)], line=1.3)
        x += w2 + gap2

    textbox(s, ML, Inches(6.58), CONTENT_W, Inches(0.4),
            [("Instrument: priced seed or SAFE. Valuation open — we would "
              "rather find the right partner than the highest cap.",
              11, CREAM, False)], line=1.25)
    notes(s, "30s. Give the range, the runway and the milestones. Do not state "
             "a valuation from the stage; keep it a conversation.")
    return s


def s14_contact(prs):
    s = new_slide(prs)
    s.shapes.add_picture(str(ASSETS / "logo_white.png"),
                         Inches(0.42), Inches(0.24), height=Inches(1.10))
    s.shapes.add_picture(str(ASSETS / "prism.png"),
                         Inches(6.90), Inches(4.05), width=Inches(6.10))

    # Split so each run is one line at this size in a 6.9in column — a third
    # wrapped line would run into the body copy below.
    textbox(s, ML, Inches(1.92), Inches(6.9), Inches(1.30),
            [("Let's make digital", 34, CREAM, True),
             ("evidence provable again.", 34, TEAL, True)], line=1.10)
    textbox(s, ML, Inches(3.24), Inches(6.3), Inches(0.9),
            [("We are looking for a lead investor and three design partners in "
              "financial services or media. A live demo takes ten minutes.",
              15, CREAM, False)], line=1.35)

    contacts = [
        ("Email", "REPLACE-WITH-REAL@EMAIL"),
        ("Repository", "github.com/muhammedsayeedurrahman/ai-authenticity-checker"),
        ("Phone", "REPLACE-WITH-REAL-NUMBER"),
    ]
    y = Inches(4.42)
    for label, value in contacts:
        textbox(s, ML, y, Inches(1.4), Inches(0.3),
                [(label.upper(), 9.5, VIOLET, True)])
        textbox(s, ML + Inches(1.42), y - Inches(0.04), Inches(5.35),
                Inches(0.34), [(value, 12.5, CREAM, False)])
        y += Inches(0.55)

    pill(s, ML, Inches(6.62), "Investor Pitch Deck", w=Inches(2.05))
    notes(s, "10s. Close on the ask, not on 'thank you'. Placeholders on this "
             "slide MUST be replaced before the deck leaves the room.")
    return s


def s15_sources(prs):
    s = new_slide(prs)
    chrome(s, "Appendix / Sources")
    heading(s, "Every number", "has a source.", size=34)
    kicker(s, "Full research file: docs/PROOFYX_COMPLETE_ANALYSIS.md in the "
              "repository.", y=1.76, size=12)

    groups = [
        ("Threat statistics", [
            "8M deepfake files by end-2025 — Bright Defense",
            "$893M AI-fraud losses, 22,364 cases — FBI IC3 2025",
            "$280K average incident — IRONSCALES, Fall 2025",
            "60% confident, 0.1% perfect — iProov, Feb 2025",
        ], MAGENTA),
        ("Market sizing", [
            "$170M to $5.6B, 47.6% CAGR — Market.us",
            "Deepfake AI market — Mordor Intelligence",
            "$15.1B by 2035 — OpenPR",
            "Defence $0.82B — MarketsandMarkets",
        ], CYAN),
        ("Regulation", [
            "EU AI Act Art. 50 / 99 — Blackbird.AI",
            "TAKE IT DOWN Act — Stack Cyber",
            "DoD investment — MIT Tech Review",
            "India IT Rules 2026 — docs/COMPLIANCE.md",
        ], VIOLET),
        ("Competition and funding", [
            "$259M across 26 companies — Tracxn",
            "Doppel $70M Series C — Fortune",
            "Reality Defender $33M — SecurityWeek",
            "GetReal $17.5M — TechCrunch",
        ], TEAL),
    ]
    gap = Inches(0.26)
    w = (CONTENT_W - 3 * gap) / 4
    x = ML
    for title, items, accent in groups:
        card(s, x, Inches(2.30), w, Inches(3.05), fill=PANEL, outline=accent,
             line_w=1.2)
        textbox(s, x + Inches(0.22), Inches(2.48), w - Inches(0.44),
                Inches(0.34), [(title.upper(), 10, accent, True)])
        runs = [(f"·  {item}", 10.5, MUTED, False) for item in items]
        textbox(s, x + Inches(0.22), Inches(2.88), w - Inches(0.44),
                Inches(2.30), runs, line=1.3, space_after=6)
        x += w + gap

    textbox(s, ML, Inches(5.62), CONTENT_W, Inches(0.6),
            [("Every chart is regenerated by scripts/make_pitch_figures.py. "
              "ProofyX accuracy and latency figures come from the CorefakeNet "
              "evaluation run on 332 held-out samples; the eval harness is in "
              "the repository and the run is repeatable on request.",
              10, MUTED, False)], line=1.25)
    notes(s, "Appendix. Do not present unless asked. Hand this slide over when "
             "an investor questions a figure.")
    return s


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    for builder in (s01_title, s02_problem, s03_why_now, s04_solution,
                    s05_product, s06_traction, s07_market, s08_opportunity,
                    s09_competition, s10_business_model, s11_financials,
                    s12_team, s13_ask, s14_contact, s15_sources):
        builder(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
